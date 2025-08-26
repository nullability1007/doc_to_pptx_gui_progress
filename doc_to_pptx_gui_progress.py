import os
import sys
import tempfile
import shutil
import threading
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk

import win32com.client as win32
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

# ==================== Poppler 路徑 ====================
def get_poppler_path():
    # PyInstaller 打包後會放在 _MEIPASS
    if getattr(sys, "_MEIPASS", False):
        return os.path.join(sys._MEIPASS, "poppler", "bin")
    else:
        # 開發時使用原始 Poppler 路徑，請改成你的路徑
        return r"C:\poppler-25.07.0\Library\bin"

POPPLER_PATH = get_poppler_path()


# ==================== 核心轉換函數 ====================
def word_to_pdf(word_path: str, pdf_path: str):
    word = win32.gencache.EnsureDispatch("Word.Application")
    word.Visible = False
    try:
        doc = word.Documents.Open(word_path)
        doc.ExportAsFixedFormat(
            OutputFileName=pdf_path,
            ExportFormat=17,
            OpenAfterExport=False,
            OptimizeFor=0,
            Range=0,
            Item=0,
            IncludeDocProps=True,
            KeepIRM=True,
            CreateBookmarks=1,
            DocStructureTags=True,
            BitmapMissingFonts=True,
            UseISO19005_1=False,
        )
        doc.Close(False)
    finally:
        word.Quit()


def build_ppt_from_images(images, out_pptx):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    total = len(images)
    for idx, img in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank_layout)
        tmpfile = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmpfile.name, "PNG")
        slide.shapes.add_picture(tmpfile.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
        tmpfile.close()
        os.unlink(tmpfile.name)
        progress_var.set(60 + int(30 * idx / total))  # PPT階段進度 60~90%
        root.update_idletasks()

    prs.save(out_pptx)


def convert_doc_to_pptx(doc_path: str, out_pptx: str, dpi=200):
    workdir = tempfile.mkdtemp(prefix="doc2pptx_")
    try:
        pdf_path = os.path.join(workdir, "temp.pdf")
        progress_var.set(10)
        root.update_idletasks()
        word_to_pdf(os.path.abspath(doc_path), pdf_path)
        progress_var.set(40)
        root.update_idletasks()

        images = convert_from_path(pdf_path, dpi=dpi, poppler_path=POPPLER_PATH)
        progress_var.set(60)
        root.update_idletasks()

        build_ppt_from_images(images, out_pptx)
        progress_var.set(100)
        root.update_idletasks()
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


# ==================== GUI ====================
def select_file():
    filepath = filedialog.askopenfilename(
        title="選擇 Word 檔案",
        filetypes=[("Word 文件", "*.docx;*.doc")]
    )
    entry_word.delete(0, tk.END)
    entry_word.insert(0, filepath)


def convert_thread():
    word_file = entry_word.get().strip()
    if not word_file or not os.path.exists(word_file):
        messagebox.showerror("錯誤", "請先選擇有效的 Word 檔案！")
        return

    base, _ = os.path.splitext(word_file)
    out_file = base + ".pptx"

    try:
        convert_doc_to_pptx(word_file, out_file)
        messagebox.showinfo("完成", f"轉換成功！\n輸出檔案：{out_file}")
    except Exception as e:
        messagebox.showerror("錯誤", f"轉換失敗：{e}")
    finally:
        progress_var.set(0)


def convert_action():
    t = threading.Thread(target=convert_thread)
    t.start()


# 建立主視窗
root = tk.Tk()
root.title("Word 轉 PPT 工具")

frame = tk.Frame(root, padx=20, pady=20)
frame.pack()

lbl = tk.Label(frame, text="Word 檔案：")
lbl.grid(row=0, column=0, sticky="w")

entry_word = tk.Entry(frame, width=40)
entry_word.grid(row=0, column=1, padx=5)

btn_browse = tk.Button(frame, text="選擇檔案", command=select_file)
btn_browse.grid(row=0, column=2, padx=5)

btn_convert = tk.Button(frame, text="轉換成 PPT", command=convert_action, bg="#4CAF50", fg="white")
btn_convert.grid(row=1, column=0, columnspan=3, pady=10)

# 進度條
progress_var = tk.IntVar()
progress_bar = ttk.Progressbar(frame, variable=progress_var, maximum=100, length=400)
progress_bar.grid(row=2, column=0, columnspan=3, pady=10)

root.mainloop()
